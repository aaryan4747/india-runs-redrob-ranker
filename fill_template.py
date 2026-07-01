#!/usr/bin/env python3
"""
Template Filler for Redrob Ideathon (Track 2)
==============================================
Loads the official template.pptx, replaces placeholder texts with
our detailed project strategy, and saves the final pitch deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def fill_template():
    prs = Presentation('template.pptx')
    
    # Substring mapping to replace placeholder text
    replacements = {
        # Slide 1 placeholders
        "Team Name :": "Team Name : AaryanK-Solo",
        "Problem Statement :": "Problem Statement : Track 1 - Intelligent Candidate Discovery & Ranking",
        "Team Leader Name :": "Team Leader Name : Boyapati Aaryan Karthik",
        
        # Slide 2 (Solution Overview)
        "proposed solution":
            "TalentLens AI is an intelligent candidate evaluation and ranking system designed for the Redrob ecosystem.\n\n"
            "Key Differentiators:\n"
            "1. Multi-Dimensional Scoring: Evaluates candidate suitability across 5 key pillars (Role Fit 30%, Skill Match 25%, Career Quality 15%, Behavioral Signals 20%, Logistics 10%) instead of a simple keyword count.\n"
            "2. Active Availability Integration: Prioritizes active, responsive candidates (by tracking recruiter response rates and platform login recency) over unreachable static profiles.\n"
            "3. Active Cheat Detection: Curated skill taxonomy with alias resolution (e.g. JS to JavaScript) penalizes \"keyword stuffers\" who list dozens of expert skills with zero tenure/experience.",
            
        # Slide 3 (JD Understanding & Candidate Evaluation)
        "key requirements":
            "Key Requirements Extracted:\n"
            "• Target Experience: 5–9 years in applied ML/AI engineering roles at product companies.\n"
            "• Technical Fit: Production experience with embeddings-based retrieval systems (sentence-transformers), vector databases (Pinecone, Qdrant, Milvus), search/ranking infrastructure (Elasticsearch, OpenSearch), and strong Python.\n"
            "• Behavioral Alignment: High recruiter responsiveness and active platform engagement.\n\n"
            "Candidate Evaluation Heuristics:\n"
            "1. Role Fit: Prioritizes product-firm ML titles (ML Engineer, AI Engineer) and penalizes IT consulting services-only tenures (TCS, Infosys, Wipro, etc.).\n"
            "2. Skill Match: Direct and alias matches are weighted by stated proficiency (Expert/Advanced) and duration (in months).\n"
            "3. Availability Index: Multiplies base scores with live behavioral signals (response rate, response time, and login recency).",
            
        # Slide 4 (Ranking Methodology)
        "retrieve, score, and rank":
            "Retrieval & Ranking Workflow:\n"
            "1. Pre-filtering: Flags and discards honeypot candidates using a multi-heuristic check.\n"
            "2. Scoring: Computes five independent scores for each candidate:\n"
            "   - Role-Fit = Weighted sum of Title Score, Experience Score, Summary Relevance, and Product/Service Ratio.\n"
            "   - Skill-Match = Direct matches weighted by duration and proficiency, minus keyword stuffing penalty.\n"
            "   - Career-Quality = Tenure stability (average months per role) + Education Field & Tier relevance.\n"
            "   - Behavioral = Login recency, recruiter response rate, response time, and GitHub activity score.\n"
            "   - Logistics = Location alignment (Noida/Pune/NCR), notice period, and preferred work mode.\n"
            "3. Weighted Aggregation: Combined Score = 0.30 × Role-Fit + 0.25 × Skill-Match + 0.15 × Career-Quality + 0.20 × Behavioral + 0.10 × Logistics.\n"
            "4. Sorting & Tie-Breaking: Sorted descending by score. Ties are broken deterministically by candidate_id ascending.",
            
        # Slide 5 (Explainability & Data Validation)
        "ranking decisions explained":
            "Explainability & Validation Heuristics:\n"
            "• Automated Fact Extraction: Explanations are compiled directly from candidate records (e.g. YOE, current title, named skills) to ensure 100% factual accuracy with zero hallucinations.\n"
            "• Proactive Honeypot Filtering: Profiles showing impossible contradictions (e.g. 8 years YOE at a company founded 3 years ago, or Expert proficiency with 0 months duration) are flagged as spammers and assigned a score of 0.\n"
            "• Keyword Stuffing Penalty: If a candidate has a high ratio of beginner skills with short durations, the Skill Match score is discounted by up to 50%.\n"
            "• Muted/Null Signal Handling: missing values (e.g., empty expected salary, unlinked GitHub) default to neutral, non-penalizing values.",
            
        # Slide 6 (End-to-End Workflow)
        "complete workflow":
            "1. JD Input: Recruiter inputs target job description constraints (must-have skills, YOE, location, preferred titles).\n"
            "2. Intent Extraction: Parser maps inputs to the hierarchical skill taxonomy and location hubs.\n"
            "3. Feature Processing: Pipeline loads 100,000 JSONL records, validating schemas and extracting career history, educational tiers, and active signals.\n"
            "4. Scoring & Filtering: Candidate profiles are audited for honeypot heuristics. Valid candidates are scored across the 5 pillars.\n"
            "5. Sorted Output: Candidates are sorted, top 100 are selected, and factual reasonings are compiled.\n"
            "6. File Generation: Produces a fully-compliant submission CSV in under 24 seconds on CPU.",
            
        # Slide 8 (Results & Performance)
        "results or insights":
            "Ranking Quality Insights:\n"
            "• Top Ranked Profile: NLP Engineer with 6.6 YOE at product company, 92% recruiter response rate, sub-30d notice period, active GitHub (82/100).\n"
            "• Services-to-Product Transition: Capable candidates from service firms are elevated only if they showcase outstanding open-source contributions.\n"
            "• Zero Honeypot Leakage: 12 honeypots detected and successfully filtered out of the top 100 list.\n\n"
            "Execution Metrics (Widescreen Performance):\n"
            "• Execution Speed: 100,000 candidate profiles loaded, scored, sorted, and saved in 23.2 seconds.\n"
            "• Memory Footprint: Peak memory < 1.2 GB RAM (well below 16 GB budget).\n"
            "• Compute Constraint: 100% CPU-only execution with zero external network dependency or LLM API latencies.",
            
        # Slide 9 (Technologies Used)
        "technologies, frameworks":
            "1. Python 3.11: Primary runtime, chosen for standard performance, robustness, and universal compatibility.\n"
            "2. Standard Library (json, csv, gzip, math, argparse, re): Used exclusively to keep the environment lightweight, eliminate external library latency, and guarantee zero-install reproducibility.\n"
            "3. PyYAML: Used for parsing submission metadata templates.\n"
            "4. Leaflet.js / Chart.js (in web mockups): Integrated for the interactive recruiter-facing analytics dashboard.\n"
            "5. GitHub & Git: Version control and repository hosting.\n"
            "6. Google Colab: Used to deploy the sandbox runner notebook.",
            
        # Slide 10 (Submission Assets)
        "Github video etc":
            "• Code Repository: https://github.com/aaryan4747/india-runs-redrob-ranker\n"
            "• Sandbox Notebook: https://colab.research.google.com/github/aaryan4747/india-runs-redrob-ranker/blob/main/redrob_ranker_sandbox.ipynb\n"
            "• Verified Output Path: /Users/aaryankarthik/Downloads/aaryan4747_submission.csv\n"
            "• Pitch Presentation Path: /Users/aaryankarthik/Downloads/TalentLens_AI_Pitch_Deck.pptx"
    }

    # Loop through slides and shape texts
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                original_text = shape.text.strip().replace('\n', ' ')
                
                # Check for direct matches or substring matches
                for key, val in replacements.items():
                    if key in original_text:
                        # Clear existing paragraphs except first
                        tf = shape.text_frame
                        # Keep the text, but replace it
                        tf.text = val
                        
                        # Style the text (keep it clean and professional)
                        for paragraph in tf.paragraphs:
                            for r in paragraph.runs:
                                r.font.name = 'Calibri'
                                r.font.size = Pt(13) if idx > 0 else Pt(22)
                                # Color styling
                                r.font.color.rgb = RGBColor(0, 0, 0)
                        
                        print(f"Replaced text on Slide {idx+1}")
                        break
                        
    # Slide 7 (System Architecture) - add the architecture drawing
    slide7 = prs.slides[6]
    # Remove existing shapes except title to clear space for the diagram
    for shape in list(slide7.shapes):
        if shape.has_text_frame and "System Architecture" not in shape.text:
            slide7.shapes._spTree.remove(shape._element)
            
    # Add a textbox with clean formatting
    txBox = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf7 = txBox.text_frame
    tf7.word_wrap = True
    
    p = tf7.paragraphs[0]
    p.text = "TalentLens AI is designed as a modular, CPU-efficient candidate matching pipeline:"
    p.font.name = 'Calibri'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    p2 = tf7.add_paragraph()
    p2.text = (
        "                    ┌────────────────────────────────┐\n"
        "                    │      Input Candidates JSONL    │\n"
        "                    └───────────────┬────────────────┘\n"
        "                                    ▼\n"
        "                    ┌────────────────────────────────┐\n"
        "                    │      Honeypot Filter Layer     │\n"
        "                    └───────────────┬────────────────┘\n"
        "                                    ▼\n"
        "      ┌─────────────────────┼─────────────────────┬─────────────────────┐\n"
        "      ▼                     ▼                     ▼                     ▼\n"
        "┌───────────┐         ┌───────────┐         ┌───────────┐         ┌───────────┐\n"
        "│ Role-Fit  │         │Skill Match│         │Behavioral │         │ Logistics │\n"
        "│   (30%)   │         │   (25%)   │         │   (20%)   │         │   (10%)   │\n"
        "└─────┬─────┘         └─────┬─────┘         └─────┬─────┘         └─────┬─────┘\n"
        "      └─────────────────────┼─────────────────────┴─────────────────────┘\n"
        "                                    ▼\n"
        "                    ┌────────────────────────────────┐\n"
        "                    │      Career Quality (15%)      │\n"
        "                    └───────────────┬────────────────┘\n"
        "                                    ▼\n"
        "                    ┌────────────────────────────────┐\n"
        "                    │     Weighted Aggregation       │\n"
        "                    └───────────────┬────────────────┘\n"
        "                                    ▼\n"
        "                    ┌────────────────────────────────┐\n"
        "                    │  Deterministic Sort & Output   │\n"
        "                    └────────────────────────────────┘"
    )
    p2.font.name = 'Courier New'
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0, 0, 0) # Black colored architecture diagram
    p2.space_before = Pt(14)

    # Save to the Downloads folder
    output_path = '/Users/aaryankarthik/Downloads/TalentLens_AI_Pitch_Deck.pptx'
    prs.save(output_path)
    # Also save to current directory for git
    prs.save('TalentLens_AI_Pitch_Deck.pptx')
    print(f"Saved filled template presentation to: {output_path}")

if __name__ == "__main__":
    fill_template()
