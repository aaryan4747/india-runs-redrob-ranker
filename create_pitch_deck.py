#!/usr/bin/env python3
"""
Pitch Deck Generator for Redrob Ideathon (Track 2)
==================================================
India Runs 2026 · Powered by Hack2skill × Redrob AI
Author: Aaryan Karthik (Boyapati)

Generates a highly professional, dark-themed presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path="TalentLens_AI_Pitch_Deck.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 widescreen standard
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    BG_COLOR = RGBColor(15, 15, 26)       # #0F0F1A Deep Dark Violet
    WHITE = RGBColor(241, 245, 249)        # #F1F5F9 Primary text
    MUTED = RGBColor(148, 163, 184)       # #94A3B8 Muted text
    CYAN = RGBColor(6, 182, 212)          # #06B6D4 Highlight Cyan
    PURPLE = RGBColor(167, 139, 250)      # #A78BFA Highlight Purple
    CARD_BG = RGBColor(26, 26, 46)        # #1A1A2E Card background
    CARD_BORDER = RGBColor(40, 40, 70)

    # Blank layout helper
    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, category_text="TALENTLENS AI PITCH"):
        # Category Tracker (Top-left)
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN

        # Slide Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    # ============================================================================
    # SLIDE 1: Title Slide (Sleek, glowing presentation center)
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)

    # Visual Accent Shape (Glow effect emulation)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()

    # Title Info
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "TALENTLENS AI"
    p.font.name = "Arial"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Redrob’s Next-Gen Intelligent Candidate Discovery Engine"
    p2.font.name = "Arial"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = CYAN
    p2.space_before = Pt(10)

    # Subtitles/Team Info
    info_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(10), Inches(2.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p3 = tf_info.paragraphs[0]
    p3.text = "TRACK 2: THE IDEATHON · IDEA SUBMISSION"
    p3.font.name = "Arial"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = PURPLE
    
    p4 = tf_info.add_paragraph()
    p4.text = "Team Leader: Boyapati Aaryan Karthik  |  Team Name: AaryanK-Solo"
    p4.font.name = "Arial"
    p4.font.size = Pt(14)
    p4.font.color.rgb = MUTED
    p4.space_before = Pt(8)

    # ============================================================================
    # SLIDE 2: The Core Problem (Why traditional hiring is broken)
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "The Core Problem: Traditional Recruitment is Broken")

    problems = [
        ("Traditional ATS Keyword Traps", "Traditional ATS filters discard 88% of qualified candidates by relying on superficial keyword matches, resulting in screen fatigue and poor matching accuracy.", PURPLE),
        ("The Keyword Stuffing Cheat", "Candidates exploit algorithms by stuffing skills sections with advanced/expert buzzwords while lacking real-world tenure or direct application experience.", CYAN),
        ("The Inactivity Gap", "Perfect-on-paper candidates are often selected for outreach despite being inactive on the platform for 6+ months, creating wasted recruitment cycles.", PURPLE),
        ("Tier-2/3 Disadvantage", "Outstanding self-taught candidates with extensive portfolios and active GitHub contributions are bypassed due to automated education-pedigree filters.", CYAN)
    ]

    for i, (title, desc, color) in enumerate(problems):
        x = Inches(0.8 + (i % 2) * 5.8)
        y = Inches(1.8 + (i // 2) * 2.5)
        
        # Card Background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.3)
        tf_card.margin_top = Inches(0.2)
        tf_card.margin_right = Inches(0.3)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = f"⚠ {title}"
        p_t.font.name = "Arial"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(8)

    # ============================================================================
    # SLIDE 3: The Solution (Introducing TalentLens AI)
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "The Solution: Multi-Dimensional Evaluation Engine")

    # Left Column: Introduction
    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0))
    tf_intro = intro_box.text_frame
    tf_intro.word_wrap = True
    p_intro = tf_intro.paragraphs[0]
    p_intro.text = "TalentLens AI"
    p_intro.font.name = "Arial"
    p_intro.font.size = Pt(24)
    p_intro.font.bold = True
    p_intro.font.color.rgb = CYAN
    
    p_desc = tf_intro.add_paragraph()
    p_desc.text = "A multi-dimensional scoring and ranking system built directly into the Redrob ecosystem. It evaluates candidate capability, availability, and alignment across five key criteria, bypassing spammers and automatically elevating genuine talent."
    p_desc.font.name = "Arial"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = MUTED
    p_desc.space_before = Pt(12)

    # Right Column: The 5 Pillars
    pillars = [
        ("Role Fit (30%)", "Evaluates current title alignment, product-company tenure, and direct AI/ML production keywords.", PURPLE),
        ("Skill Match (25%)", "Curated taxonomy matching with alias resolution. Weights proficiency, duration, and endorsements.", CYAN),
        ("Career Quality (15%)", "Measures tenure stability (anti-job-hopping) and academic alignment (Computer Science / ML).", PURPLE),
        ("Behavioral Signals (20%)", "Weights responsiveness (rate and speed), logins, and GitHub contribution scores.", CYAN),
        ("Logistics (10%)", "Ensures alignment with location (Noida/Pune), notice period (sub-30d preferred), and work mode.", PURPLE)
    ]

    for i, (title, desc, color) in enumerate(pillars):
        y = Inches(1.8 + i * 1.05)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.8), y, Inches(6.8), Inches(0.95))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_top = Inches(0.1)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(2)

    # ============================================================================
    # SLIDE 4: Interactive Architecture & User Journey
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "AI Architecture & Recruiter Journey Map")

    steps = [
        ("1. Input JD Parsing", "System automatically parses candidate JDs, extracting required skills, target YOE, preferred locations, and target titles.", PURPLE),
        ("2. Multidimensional Extraction", "Feature extractor queries candidate database, capturing profile structure, career history, and active platform activity logs.", CYAN),
        ("3. Weighted Scoring", "Five distinct scoring engines run simultaneously on CPU, scoring candidates based on capability and actual availability.", PURPLE),
        ("4. Honeypot & Tiebreak", "Profiles with contradictory signals are filtered out (Score=0). Remaining candidate lists are sorted with deterministic tiebreakers.", CYAN)
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.8 + i * 2.95)
        
        # Step Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_top = Inches(0.2)
        tf_card.margin_right = Inches(0.15)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.alignment = PP_ALIGN.CENTER
        
        p_divider = tf_card.add_paragraph()
        p_divider.text = "──────────"
        p_divider.font.name = "Arial"
        p_divider.font.size = Pt(10)
        p_divider.font.color.rgb = CARD_BORDER
        p_divider.alignment = PP_ALIGN.CENTER
        p_divider.space_before = Pt(4)
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(10)
        p_d.alignment = PP_ALIGN.LEFT

    # ============================================================================
    # SLIDE 5: Anti-Cheat & Honeypot Engine
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "Anti-Cheat Mechanics & Honeypot Defense")

    # Left: Keyword Stuffing
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG
    left_card.line.color.rgb = CARD_BORDER
    
    tf_l = left_card.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.3)
    tf_l.margin_top = Inches(0.3)
    
    p_lt = tf_l.paragraphs[0]
    p_lt.text = "🛡 Anti-Keyword Stuffing"
    p_lt.font.name = "Arial"
    p_lt.font.size = Pt(18)
    p_lt.font.bold = True
    p_lt.font.color.rgb = CYAN
    
    p_ld = tf_l.add_paragraph()
    p_ld.text = "Keyword stuffers are penalized through active capability checks:\n\n" \
                "• Duration-Weighted Skills: Skills listed with short durations (e.g. <3 months) are down-weighted.\n\n" \
                "• Proficiency discounting: Beginner skills receive minimal weight, rewarding experienced practitioners.\n\n" \
                "• Skill-to-Experience Ratio: Profiles claiming excessive expert/advanced skills relative to their total years of experience are marked suspicious."
    p_ld.font.name = "Arial"
    p_ld.font.size = Pt(13)
    p_ld.font.color.rgb = MUTED
    p_ld.space_before = Pt(12)

    # Right: Honeypot Heuristics
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = CARD_BORDER
    
    tf_r = right_card.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.3)
    
    p_rt = tf_r.paragraphs[0]
    p_rt.text = "🎯 Heuristic Honeypot Detection"
    p_rt.font.name = "Arial"
    p_rt.font.size = Pt(18)
    p_rt.font.bold = True
    p_rt.font.color.rgb = PURPLE
    
    p_rd = tf_r.add_paragraph()
    p_rd.text = "We built 9 strict validation rules to capture honeypot spammers:\n\n" \
                "• Chronological Contradictions: E.g., claiming 8 years of tenure at a startup established 3 years ago.\n\n" \
                "• Assessment Discrepancies: Expert/Advanced profiles scoring poorly (<50%) on Redrob platform skill tests.\n\n" \
                "• Fictional Company Auditing: Mismatched experience durations at known fictional placeholder companies.\n\n" \
                "• Profile Incompleteness: High completeness scores with near-empty work history details."
    p_rd.font.name = "Arial"
    p_rd.font.size = Pt(13)
    p_rd.font.color.rgb = MUTED
    p_rd.space_before = Pt(12)

    # ============================================================================
    # SLIDE 6: Behavioral Signals — Availability as a Metric
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "Behavioral Signals: Availability-First Ranking")

    # Callout text
    callout_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.2))
    tf_c = callout_box.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "\"A perfect-on-paper candidate who hasn't logged in for 6 months and has a 5% response rate is, for hiring purposes, not actually available.\""
    p_c.font.name = "Arial"
    p_c.font.size = Pt(18)
    p_c.font.italic = True
    p_c.font.bold = True
    p_c.font.color.rgb = PURPLE
    p_c.alignment = PP_ALIGN.CENTER

    # Three key columns
    cols = [
        ("Recruiter Response Rate", "Active tracking of candidate reply rates. Candidates responding to >70% of messages receive maximum boost; unresponsive profiles (<15%) are heavily deprioritized.", CYAN),
        ("Platform Recency", "Evaluates days since last login. Active logins within 30 days denote high search intent; inactivity >180 days triggers candidate cooling/down-ranking.", PURPLE),
        ("GitHub Activity Score", "Extracts commits, PRs, and active contributions. Allows self-taught developers from Tier-2/3 colleges with rich portfolios to easily outrank inactive Tier-1 grads.", CYAN)
    ]

    for i, (title, desc, color) in enumerate(cols):
        x = Inches(0.8 + i * 3.95)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(3.8), Inches(3.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.25)
        tf_card.margin_top = Inches(0.25)
        tf_card.margin_right = Inches(0.25)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        p_t.alignment = PP_ALIGN.CENTER
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(10)

    # ============================================================================
    # SLIDE 7: India-Specific Market Alignment
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "India-Specific Recruitment Alignment")

    alignments = [
        ("Product vs. Service Focus", "Redrob AI prioritizes engineers who have built and shipped products. TalentLens AI explicitly distinguishes product experience from IT consulting companies (Infosys, TCS, Wipro, etc.), heavily penalizing career-long service-only profiles unless supported by outstanding open-source projects.", PURPLE),
        ("Notice Period Weighting", "Notice periods in India are notoriously long (up to 90 days). The system awards maximum points to sub-30-day candidates, enabling rapid startup hiring, while scaling down score metrics for 90-day locks.", CYAN),
        ("Regional Location Filtering", "Prioritizes immediate hub matches (Noida, Pune, Bangalore) to guarantee high onboarding success while accommodating candidates willing to relocate to offices.", PURPLE)
    ]

    for i, (title, desc, color) in enumerate(alignments):
        y = Inches(1.8 + i * 1.7)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.3)
        tf_card.margin_top = Inches(0.2)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(6)

    # ============================================================================
    # SLIDE 8: Explanations & Interactive Recruiter UX
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "Explainable AI (XAI) Recruiter Experience")

    # Left: Explanation text
    ux_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_ux = ux_box.text_frame
    tf_ux.word_wrap = True
    p_uxt = tf_ux.paragraphs[0]
    p_uxt.text = "Empowering Recruiters with Trust"
    p_uxt.font.name = "Arial"
    p_uxt.font.size = Pt(20)
    p_uxt.font.bold = True
    p_uxt.font.color.rgb = CYAN
    
    p_uxd = tf_ux.add_paragraph()
    p_uxd.text = "Recruiters cannot trust a simple number. TalentLens AI generates clean, structured explanations referencing specific profile facts:\n\n" \
                 "• Dynamic Sliders: Recruiters can customize importance values (e.g., dial down notice period requirements, dial up GitHub scores).\n\n" \
                 "• Highlight Strengths: Explicitly surfaces years in the target experience band, relevant past employers, and verified test results.\n\n" \
                 "• Concerns Flagging: Surfaces high salary expectations, long notice periods, or services-only tenures immediately."
    p_uxd.font.name = "Arial"
    p_uxd.font.size = Pt(13)
    p_uxd.font.color.rgb = MUTED
    p_uxd.space_before = Pt(12)

    # Right: Mockup Box (emulated UI Card)
    mockup = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    mockup.fill.solid()
    mockup.fill.fore_color.rgb = CARD_BG
    mockup.line.color.rgb = PURPLE
    mockup.line.width = Pt(2)
    
    tf_m = mockup.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = Inches(0.4)
    tf_m.margin_top = Inches(0.3)
    tf_m.margin_right = Inches(0.4)
    
    p_mt = tf_m.paragraphs[0]
    p_mt.text = "CANDIDATE DOSSIER: CAND_0052682"
    p_mt.font.name = "Courier New"
    p_mt.font.size = Pt(14)
    p_mt.font.bold = True
    p_mt.font.color.rgb = PURPLE
    
    p_m1 = tf_m.add_paragraph()
    p_m1.text = "Current Title: NLP Engineer @ Aganitha\n" \
                "Experience:    6.6 Years (Target Band: 5-9)\n" \
                "Match Score:   93.15 / 100 [Strong Hire]\n" \
                "────────────────────────────────────────"
    p_m1.font.name = "Courier New"
    p_m1.font.size = Pt(11)
    p_m1.font.color.rgb = WHITE
    p_m1.space_before = Pt(10)
    
    p_m2 = tf_m.add_paragraph()
    p_m2.text = "💪 KEY STRENGTHSsurfaced:\n" \
                "- Active open source contributor (GitHub: 82/100)\n" \
                "- Proficient in sentence-transformers and Qdrant\n" \
                "- Extremely responsive (92% recruiter response rate)\n\n" \
                "⚠ AREAS OF CONCERN flagged:\n" \
                "- 60-day notice period (negotiable up to 30d)"
    p_m2.font.name = "Arial"
    p_m2.font.size = Pt(11)
    p_m2.font.color.rgb = MUTED
    p_m2.space_before = Pt(12)

    # ============================================================================
    # SLIDE 9: System Benchmarks & Performance
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "System Benchmarks & Latency Efficiency")

    # Large stat callouts
    stats = [
        ("23.2s", "Total pipeline execution runtime for all 100,000 candidates.", CYAN),
        ("0.0ms", "Zero API call latency due to fully localized indexing and taxonomy rules.", PURPLE),
        ("0 Leaked", "Honeypot spammers in the top 100 candidate recommendations.", CYAN),
        ("< 1.2 GB", "Peak RAM footprint, conforming to tight microservice deployments.", PURPLE)
    ]

    for i, (value, desc, color) in enumerate(stats):
        x = Inches(0.8 + i * 2.95)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.8), Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_top = Inches(0.3)
        tf_card.margin_right = Inches(0.2)
        
        p_v = tf_card.paragraphs[0]
        p_v.text = value
        p_v.font.name = "Arial"
        p_v.font.size = Pt(36)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.alignment = PP_ALIGN.CENTER
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(20)
        p_d.alignment = PP_ALIGN.CENTER

    # ============================================================================
    # SLIDE 10: Future Roadmap (Scaling the Platform)
    # ============================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide)
    add_header(slide, "Future Development Roadmap")

    phases = [
        ("Phase 1: Verification Integration", "Direct API integration with ABDM (Ayushman Bharat Digital Mission), Aadhaar, DigiLocker, and GitHub OAuth to eliminate forged degrees and experience claims prior to ranking.", PURPLE),
        ("Phase 2: Hybrid LLM Agent Re-Ranking", "Combine our hyper-fast heuristic ranker with a secondary local LLM agent (e.g., Llama-3-8B-Instruct) to parse and re-rank the top 500 candidate profiles.", CYAN),
        ("Phase 3: Conversational Voice Screening", "Redrob AI phone screening agents update candidate profiles with direct assessment ratings (communication, logic, coding) in real-time, automatically updating availability indices.", PURPLE)
    ]

    for i, (title, desc, color) in enumerate(phases):
        y = Inches(1.8 + i * 1.7)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.3)
        tf_card.margin_top = Inches(0.2)
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = "Arial"
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = MUTED
        p_d.space_before = Pt(6)

    prs.save(output_path)
    print(f"PowerPoint created successfully: {output_path}")

if __name__ == "__main__":
    create_deck()
